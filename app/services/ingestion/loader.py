"""
Document loader service supporting PDF, TXT, and Markdown files.
Uses PyMuPDF (fitz) for PDF extraction.
"""

import os
import logging
import re
from typing import List, Optional, Iterator
import fitz  # PyMuPDF
from app.services.ingestion.schemas import LoadedDocument, LoadedPage, StreamedPage

logger = logging.getLogger(__name__)


class DocumentLoader:
    """Service for loading files and extracting text with metadata via streaming."""

    def __init__(self) -> None:
        logger.info("Initialized DocumentLoader")

    def detect_file_type(self, file_path: str) -> str:
        """Determines document type based on file extension."""
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".pdf":
            return "pdf"
        elif ext in [".md", ".markdown"]:
            return "md"
        elif ext == ".txt":
            return "txt"
        elif ext == ".docx":
            return "docx"
        elif ext == ".pptx":
            return "pptx"
        else:
            raise ValueError(f"Unsupported file extension: {ext}. Supported: pdf, txt, md, docx, pptx.")

    def stream_file(self, file_path: str) -> Iterator[StreamedPage]:
        """Streams a document from file_path, yielding StreamedPage objects."""
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        doc_type = self.detect_file_type(file_path)
        source_name = os.path.basename(file_path)
        logger.info(f"Streaming {doc_type.upper()} document: {source_name}")

        if doc_type == "pdf":
            yield from self._stream_pdf(file_path, source_name, doc_type)
        elif doc_type == "md":
            yield from self._stream_markdown(file_path, source_name, doc_type)
        elif doc_type == "docx":
            yield from self._stream_docx(file_path, source_name, doc_type)
        elif doc_type == "pptx":
            yield from self._stream_pptx(file_path, source_name, doc_type)
        else:
            yield from self._stream_txt(file_path, source_name, doc_type)

    def load_file(self, file_path: str) -> LoadedDocument:
        """Legacy compatibility wrapper. Avoid for large documents."""
        doc_type = self.detect_file_type(file_path)
        source_name = os.path.basename(file_path)
        pages = [sp.page for sp in self.stream_file(file_path)]
        return LoadedDocument(source=source_name, document_type=doc_type, pages=pages)

    def _stream_pdf(self, file_path: str, source: str, doc_type: str) -> Iterator[StreamedPage]:
        try:
            doc = fitz.open(file_path)
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text = page.get_text("text").strip()
                if not text:
                    continue
                
                section_title = None
                lines = [line.strip() for line in text.splitlines() if line.strip()]
                if lines:
                    first_line = lines[0]
                    if len(first_line) < 60 and (first_line.isupper() or first_line.istitle()):
                        section_title = first_line

                loaded_page = LoadedPage(text=text, page_number=page_num + 1, section_title=section_title)
                yield StreamedPage(source=source, document_type=doc_type, page=loaded_page)
            doc.close()
        except Exception as e:
            logger.error(f"Error streaming PDF {file_path}: {e}")
            raise RuntimeError(f"Failed to stream PDF: {e}")

    def _stream_markdown(self, file_path: str, source: str, doc_type: str) -> Iterator[StreamedPage]:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()

            sections = re.split(r"(^#{1,2}\s+.*)", content, flags=re.MULTILINE)
            current_title = None
            buffer = []

            for part in sections:
                if not part.strip():
                    continue
                match = re.match(r"^#{1,2}\s+(.*)", part.strip())
                if match:
                    if buffer:
                        loaded_page = LoadedPage(text="\n".join(buffer).strip(), section_title=current_title)
                        yield StreamedPage(source=source, document_type=doc_type, page=loaded_page)
                        buffer = []
                    current_title = match.group(1).strip()
                else:
                    buffer.append(part.strip())

            if buffer:
                loaded_page = LoadedPage(text="\n".join(buffer).strip(), section_title=current_title)
                yield StreamedPage(source=source, document_type=doc_type, page=loaded_page)

        except Exception as e:
            logger.error(f"Error streaming Markdown {file_path}: {e}")
            raise RuntimeError(f"Failed to stream Markdown: {e}")

    def _stream_txt(self, file_path: str, source: str, doc_type: str) -> Iterator[StreamedPage]:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read().strip()
            if text:
                yield StreamedPage(source=source, document_type=doc_type, page=LoadedPage(text=text))
        except Exception as e:
            logger.error(f"Error streaming TXT {file_path}: {e}")
            raise RuntimeError(f"Failed to stream TXT: {e}")

    def _stream_docx(self, file_path: str, source: str, doc_type: str) -> Iterator[StreamedPage]:
        try:
            from docx import Document
            doc = Document(file_path)
            buffer = []
            page_num = 1
            section_title = None
            
            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                    
                if para.style.name.startswith('Heading'):
                    if buffer:
                        loaded_page = LoadedPage(text="\n".join(buffer), page_number=page_num, section_title=section_title)
                        yield StreamedPage(source=source, document_type=doc_type, page=loaded_page)
                        buffer = []
                        page_num += 1
                    section_title = text
                else:
                    buffer.append(text)
                    
            if buffer:
                loaded_page = LoadedPage(text="\n".join(buffer), page_number=page_num, section_title=section_title)
                yield StreamedPage(source=source, document_type=doc_type, page=loaded_page)
                
        except ImportError:
            logger.error("python-docx is not installed. Please run: pip install python-docx")
            raise
        except Exception as e:
            logger.error(f"Error streaming DOCX {file_path}: {e}")
            raise RuntimeError(f"Failed to stream DOCX: {e}")

    def _stream_pptx(self, file_path: str, source: str, doc_type: str) -> Iterator[StreamedPage]:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            
            for i, slide in enumerate(prs.slides):
                buffer = []
                section_title = None
                
                if slide.shapes.title and slide.shapes.title.has_text_frame:
                    section_title = slide.shapes.title.text.strip()
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        if shape == slide.shapes.title:
                            continue
                        buffer.append(shape.text.strip())
                        
                if buffer or section_title:
                    text = "\n".join(buffer) if buffer else section_title
                    loaded_page = LoadedPage(text=text, page_number=i+1, section_title=section_title)
                    yield StreamedPage(source=source, document_type=doc_type, page=loaded_page)
                    
        except ImportError:
            logger.error("python-pptx is not installed. Please run: pip install python-pptx")
            raise
        except Exception as e:
            logger.error(f"Error streaming PPTX {file_path}: {e}")
            raise RuntimeError(f"Failed to stream PPTX: {e}")
